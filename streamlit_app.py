import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import io
import textwrap
import re

st.title("PPT 자동 생성기 (Streamlit)")

st.markdown('''
내용을 단락락별로 번호로 구분해서 정리해서 넣기 (예: 1. 사업의 필요성 2. 사업의 개요 ... 5. 기대효과 형태로)
''')

st.markdown("""
**내용 입력**  
""")

# 1. 입력
content = st.text_area("PPT로 만들고 싶은 내용을 입력하세요...", height=200)
page_count = st.number_input("페이지 수", min_value=1, value=5, step=1)

# 2. 슬라이드별로 분할
items = re.findall(r'(\d+)\.\s*([^\n]+)\n([^\n]+(?:\n(?!\d+\.).+)*)', content, re.MULTILINE)
slides_content = []
for idx, title, item_content in items:
    slides_content.append({
        'title': title.strip(),
        'content': item_content.strip()
    })

# 3. 슬라이드별로 수정 UI 제공
st.markdown("**슬라이드별 내용 미리보기 및 수정**")
edited_slides = []
for i, slide in enumerate(slides_content):
    title = st.text_input(f"슬라이드 {i+1} 제목", value=slide['title'], key=f"title_{i}")
    body = st.text_area(f"슬라이드 {i+1} 내용", value=slide['content'], key=f"body_{i}")
    edited_slides.append({'title': title, 'content': body})

# 4. PPT 생성
if st.button("PPT 생성"):
    # 페이지 수에 맞게 균등 분할
    def group_sections_by_page(sections, page_count):
        n = len(sections)
        base = n // page_count
        extra = n % page_count
        grouped = []
        idx = 0
        for i in range(page_count):
            count = base + (1 if i < extra else 0)
            grouped.append(sections[idx:idx+count])
            idx += count
        return grouped

    grouped_slides = group_sections_by_page(edited_slides, int(page_count))

    prs = Presentation()
    for group in grouped_slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        y_offset = 0.5
        for slide_content in group:
            # 제목
            title_shape = slide.shapes.add_textbox(Inches(0.7), Inches(y_offset), Inches(6.5), Inches(1))
            title_frame = title_shape.text_frame
            title_frame.clear()
            p_title = title_frame.add_paragraph()
            p_title.text = slide_content.get('title', '')
            p_title.font.size = Pt(36)
            p_title.font.bold = True
            p_title.alignment = PP_ALIGN.LEFT  # 왼쪽 정렬
            y_offset += 1
            # 내용
            content_shape = slide.shapes.add_textbox(Inches(0.7), Inches(y_offset), Inches(6.5), Inches(3.5))
            content_frame = content_shape.text_frame
            content_frame.clear()
            for line in slide_content.get('content', '').split('\n'):
                if line.strip():
                    for wrapped_line in textwrap.wrap(line.strip(), width=40):
                        p = content_frame.add_paragraph()
                        p.text = wrapped_line
                        p.font.size = Pt(20)
                        p.alignment = PP_ALIGN.LEFT  # 왼쪽 정렬
            y_offset += 2.5
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    st.download_button("PPT 다운로드", ppt_io.getvalue(), file_name="output.pptx")
